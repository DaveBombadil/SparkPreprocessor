import sys
import subprocess
import importlib.util

def check_install_package(package_name: str, import_name: str = None):
    """Check if a package is installed, if not install it"""
    if import_name is None:
        import_name = package_name
    
    if importlib.util.find_spec(import_name) is None:
        print(f"Installing {package_name}...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])
            print(f"Successfully installed {package_name}")
        except subprocess.CalledProcessError as e:
            print(f"Error installing {package_name}: {str(e)}")
            sys.exit(1)
    else:
        print(f"{package_name} is already installed")

def create_spark_preprocessor_presentation(output_path: str):
    """Create the presentation with the given output path"""
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SparkPreprocessor"
    subtitle.text = "A Comprehensive Data Analysis and Preprocessing Tool for PySpark"

    # Overview Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Overview"
    tf = body.text_frame
    
    p = tf.add_paragraph()
    p.text = "SparkPreprocessor is a high-performance tool that combines:"
    
    for point in [
        "Automated data type detection and validation",
        "Comprehensive data analysis and statistics",
        "Configurable preprocessing pipeline",
        "Interactive visualizations",
        "Efficient caching system",
        "Progress tracking and logging"
    ]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 1

    # Key Features Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Key Features"
    tf = body.text_frame
    
    features = {
        "Intelligent Type Detection": [
            "Numeric/Categorical/Text/DateTime/Binary",
            "Custom type overrides",
            "Complex type handling (Arrays, Maps, Structs)"
        ],
        "Statistical Analysis": [
            "Basic statistics and distributions",
            "Missing value analysis",
            "Correlation analysis (Pearson/Chatterjee)",
            "Cardinality analysis"
        ],
        "Data Quality": [
            "Null value detection and handling",
            "Outlier detection and treatment",
            "Data type validation",
            "Duplicate handling"
        ]
    }
    
    for feature, points in features.items():
        p = tf.add_paragraph()
        p.text = feature
        
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.level = 1

    # Configuration Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Configuration Options"
    tf = body.text_frame
    
    configs = {
        "Data Preprocessing": [
            "scaling_method: standard/minmax/none",
            "null_strategy: mean/median/mode/constant/drop/flag",
            "outlier_strategy: remove/cap/scale",
            "correlation_method: pearson/chatterjee"
        ],
        "Advanced Settings": [
            "Cardinality control (max/min thresholds)",
            "Text processing parameters",
            "Column exclusions via no_process",
            "Custom type overrides"
        ]
    }
    
    for section, points in configs.items():
        p = tf.add_paragraph()
        p.text = section
        
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.level = 1

    # Usage Example Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Usage Example"
    tf = body.text_frame
    
    code = """
        # Configure preprocessing
        config = PreprocessingConfig(
            scaling_method='standard',
            null_strategy='mean',
            outlier_strategy='cap',
            correlation_method='pearson'
        )

        # Define column type overrides
        column_overrides = {
            'numeric': ['amount', 'score'],
            'text': ['description', 'comments'],
            'categorical': ['status', 'category'],
            'ignore': ['id', 'created_at']
        }

        # Initialize and run
        preprocessor = SparkPreprocessor(
            df=spark_df,
            config=config,
            column_type_overrides=column_overrides
        )

        # Get processed dataframe
        processed_df = preprocessor.fit_transform(spark_df)
    """
    
    for line in code.strip().split('\n'):
        p = tf.add_paragraph()
        p.text = line

    # Visualization Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Built-in Visualizations"
    tf = body.text_frame
    
    visualizations = [
        "Distribution plots for numeric columns",
        "Missing value heatmaps",
        "Correlation matrices",
        "Categorical value distributions",
        "Interactive profiling reports"
    ]
    
    p = tf.add_paragraph()
    p.text = "Includes various visualization methods:"
    
    for viz in visualizations:
        p = tf.add_paragraph()
        p.text = f"• {viz}"
        p.level = 1

    # Benefits Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Benefits"
    tf = body.text_frame
    
    benefits = [
        "Reduces data preprocessing time by 80%",
        "Ensures consistent data quality standards",
        "Provides comprehensive data insights",
        "Handles large-scale data efficiently",
        "Highly configurable and extensible",
        "Built-in error handling and logging",
        "Memory-efficient caching system"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = f"• {benefit}"

    # Save presentation
    prs.save(output_path)

def main():
    """Main function to create presentation with automatic dependency management"""
    print("Checking and installing required packages...")
    
    # Check and install required packages
    packages = {
        "python-pptx": "pptx",  # package_name: import_name
    }
    
    for package_name, import_name in packages.items():
        check_install_package(package_name, import_name)
    
    # Now import the required packages
    global Presentation, Inches, Pt, RGBColor, PP_ALIGN
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    print("\nCreating presentation...")
    create_spark_preprocessor_presentation("spark_preprocessor_pitch.pptx")

if __name__ == "__main__":
    try:
        print("Starting presentation generator...")
        print("This script will check for required packages and install them if necessary.")
        print("You might see some installation messages...\n")
        
        main()
        
        print("\nPresentation has been created successfully!")
        print("You should see a file named 'spark_preprocessor_pitch.pptx' in the same folder as this script.")
        print("\nNote: If you don't see the file, check the current working directory by running:")
        print("import os; print(os.getcwd())")
        
        # Keep console window open on Windows
        if sys.platform.startswith('win'):
            input("\nPress Enter to exit...")
            
    except Exception as e:
        print(f"\nAn error occurred: {str(e)}")
        print("\nPlease try the following:")
        print("1. Make sure you have Python installed correctly")
        print("2. Try running 'pip install python-pptx' manually")
        print("3. Check if you have write permissions in the current folder")
        print("\nIf the problem persists, please contact for support")
        
        # Keep console window open on Windows
        if sys.platform.startswith('win'):
            input("\nPress Enter to exit...")